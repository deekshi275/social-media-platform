from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "HostelGram: Social Media App for Hostel Students"
subtitle.text = "A Django-based platform for sharing posts, stories, and messages"

# Slide 2: Project Aim
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Project Aim"
content = slide.placeholders[1]
content.text = "To develop a dedicated social media platform for hostel students that enables easy sharing of posts, stories, and direct messaging, fostering community interaction and communication among residents."

# Slide 3: Scope
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Scope"
content = slide.placeholders[1]
content.text = "- User registration and authentication\n- Image posting with captions\n- Stories with 12-hour auto-delete\n- Real-time messaging system\n- Profile management and editing\n- Like and comment features\n- Follow/unfollow functionality\n- Search and explore users/posts"

# Slide 4: Problem Statement
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Problem Statement"
content = slide.placeholders[1]
content.text = "Hostel students lack a dedicated platform for social interaction within their community. Existing social media platforms are general-purpose and do not cater to the specific needs of hostel residents for sharing daily life, events, and maintaining connections. There is a need for a secure, easy-to-use app focused on hostel communities."

# Slide 5: Management
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Project Management"
content = slide.placeholders[1]
content.text = "- Framework: Django 6.0.5\n- Database: SQLite\n- Frontend: HTML, CSS, JavaScript\n- Media Handling: File uploads for images\n- Security: User authentication and session management\n- Deployment: Local development server"

# Slide 6: Snapshots
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Application Snapshots"
content = slide.placeholders[1]
content.text = "- Main Feed: Displays posts from followed users and stories\n- Profile Page: Shows user bio, posts, followers count\n- Messaging: Direct chat interface with message history\n- Upload Modal: Multi-image post creation\n- Story Viewer: Grouped stories from users\n- Explore Page: Browse all posts and users"

# Save the presentation
prs.save('HostelGram_Presentation.pptx')
print("Presentation created successfully!")